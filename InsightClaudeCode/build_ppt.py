"""
Claude Code 精密拆解 — PPT 生成脚本
字体规范：中文楷体，英文 Times New Roman
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree

# ── 字体名 ────────────────────────────────────────────────
ZH = "楷体"
EN = "Times New Roman"

# ── 颜色 ─────────────────────────────────────────────────
NAVY   = RGBColor(0x0f, 0x17, 0x2a)
BLUE   = RGBColor(0x1d, 0x4e, 0xd8)
LBLUE  = RGBColor(0x3b, 0x82, 0xf6)
SKYBLUE= RGBColor(0xdb, 0xea, 0xfe)
ACCENT = RGBColor(0xe0, 0xf2, 0xfe)
WHITE  = RGBColor(0xff, 0xff, 0xff)
DARK   = RGBColor(0x1e, 0x29, 0x3b)
GRAY   = RGBColor(0x64, 0x74, 0x8b)
LGRAY  = RGBColor(0xf1, 0xf5, 0xf9)
GREEN  = RGBColor(0x05, 0x96, 0x69)
LGREEN = RGBColor(0xd1, 0xfa, 0xe5)
ORANGE = RGBColor(0xd9, 0x77, 0x06)
LORANG = RGBColor(0xfe, 0xf3, 0xc7)
PURPLE = RGBColor(0x6d, 0x28, 0xd9)
LPURPL = RGBColor(0xed, 0xe9, 0xfe)
RED    = RGBColor(0xdc, 0x26, 0x26)
LRED   = RGBColor(0xfe, 0xe2, 0xe2)
TEAL   = RGBColor(0x04, 0x7a, 0x85)
LTEAL  = RGBColor(0xcf, 0xfa, 0xf8)
PINK   = RGBColor(0xa2, 0x1c, 0xaf)
LPINK  = RGBColor(0xfd, 0xe8, 0xff)

# ── 核心辅助函数 ──────────────────────────────────────────

def _set_font(run, name):
    """同时设置 latin 和 eastAsia 字体"""
    rPr = run._r.get_or_add_rPr()
    for tag in (qn('a:latin'), qn('a:ea')):
        el = rPr.find(tag)
        if el is None:
            el = etree.SubElement(rPr, tag)
        el.set('typeface', name)

def is_cjk(text):
    return any('\u4e00' <= c <= '\u9fff' for c in text)

def run_add(para, text, pt, bold=False, italic=False, color=DARK, font=None):
    r = para.add_run()
    r.text = text
    r.font.size = Pt(pt)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    fname = font or (ZH if is_cjk(text) else EN)
    r.font.name = fname
    _set_font(r, fname)
    return r

def tf_para(tf, text, pt, bold=False, italic=False, color=DARK,
            align=PP_ALIGN.LEFT, font=None, spc_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if spc_before:
        pPr = p._p.get_or_add_pPr()
        pPr.set('spcBef', str(spc_before * 100))
    run_add(p, text, pt, bold=bold, italic=italic, color=color, font=font)
    return p

def add_txt(sld, x, y, w, h, text, pt, bold=False, color=DARK,
            align=PP_ALIGN.LEFT, font=None, wrap=True):
    tb = sld.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run_add(p, text, pt, bold=bold, color=color, font=font)
    return tb

def add_rect(sld, x, y, w, h, fill=SKYBLUE, line_color=None,
             rounded=False, line_pt=1.0):
    from pptx.util import Emu
    shape_type = (MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
                  if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    shp = sld.shapes.add_shape(shape_type,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        shp.adjustments[0] = 0.04
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_pt)
    else:
        shp.line.fill.background()
    return shp

def add_box(sld, x, y, w, h, text, pt, bold=False,
            fill=SKYBLUE, text_color=DARK, line_color=None,
            align=PP_ALIGN.CENTER, font=None, rounded=True):
    shp = add_rect(sld, x, y, w, h, fill=fill, line_color=line_color, rounded=rounded)
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run_add(p, text, pt, bold=bold, color=text_color, font=font)
    return shp

def add_arrow_down(sld, x, y, w=0.18, h=0.28, color=BLUE):
    shp = sld.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def add_arrow_right(sld, x, y, w=0.35, h=0.2, color=BLUE):
    shp = sld.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def slide_header(sld, title, subtitle=""):
    add_rect(sld, 0, 0, 13.33, 1.05, fill=NAVY, rounded=False)
    add_txt(sld, 0.5, 0.18, 12.5, 0.65, title, 24, bold=True, color=WHITE)
    if subtitle:
        add_txt(sld, 0.5, 0.72, 12.5, 0.3, subtitle, 11, color=RGBColor(0xba,0xd0,0xf8))

# ══════════════════════════════════════════════════════════
# 创建 Presentation
# ══════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

W, H = 13.33, 7.5

# ══════════════════════════════════════════════════════════
# 幻灯片 01 — 封面
# ══════════════════════════════════════════════════════════
sld = prs.slides.add_slide(blank)
add_rect(sld, 0, 0, W, H, fill=NAVY, rounded=False)
add_rect(sld, 0, 0, 5.0, H, fill=RGBColor(0x07,0x0f,0x1e), rounded=False)
# 右侧竖条装饰
for i, alpha in enumerate([LBLUE, RGBColor(0x1e,0x40,0x8a), RGBColor(0x10,0x24,0x52)]):
    add_rect(sld, W - 0.15*(i+1) - 0.06, 0, 0.06, H, fill=alpha, rounded=False)

# 主标题
tb = sld.shapes.add_textbox(Inches(0.65), Inches(1.5), Inches(9), Inches(1.5))
tf = tb.text_frame
tf.word_wrap = False
p1 = tf.paragraphs[0]
p1.alignment = PP_ALIGN.LEFT
run_add(p1, "Claude Code", 50, bold=True, color=WHITE, font=EN)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.LEFT
run_add(p2, "精密拆解", 50, bold=True, color=RGBColor(0x93,0xc5,0xfd), font=ZH)

add_txt(sld, 0.65, 3.3, 10, 0.55,
        "一台精密的 Harness 工程机器  ·  七大子系统全解析", 17,
        color=RGBColor(0x7d,0xa8,0xe8), font=ZH)

# 标签
tags = [("四层架构",BLUE), ("QueryLoop",GREEN), ("权限引擎",ORANGE),
        ("工具并发",PURPLE), ("记忆体系",TEAL), ("MCP·Hooks·Skills",RED)]
tx = 0.65
for tag, c in tags:
    w_t = len(tag) * 0.21 + 0.4
    add_box(sld, tx, 4.15, w_t, 0.38, tag, 12, bold=True,
            fill=RGBColor(0x0d,0x25,0x55), text_color=RGBColor(0x93,0xc5,0xfd),
            rounded=True)
    tx += w_t + 0.18

add_txt(sld, 0.65, 6.75, 11, 0.45,
        "基于源码 zhangbo2008/claude_code_annotated + shareAI-lab/learn-claude-code",
        11, color=RGBColor(0x3a,0x45,0x5e), font=EN)

# 幻灯片编号
add_txt(sld, 12.7, 0.3, 0.5, 0.3, "01", 11, color=RGBColor(0x3a,0x45,0x5e))
print("✓ 01 封面")

# ══════════════════════════════════════════════════════════
# 幻灯片 02 — 目录
# ══════════════════════════════════════════════════════════
sld = prs.slides.add_slide(blank)
add_rect(sld, 0, 0, W, H, fill=LGRAY, rounded=False)
add_rect(sld, 0, 0, 0.1, H, fill=BLUE, rounded=False)
slide_header(sld, "目  录")

items = [
    ("01","这台机器是什么？",    "Harness Engineering Paradigm"),
    ("02","四层体系总览",         "L1 Execution → L4 Extension"),
    ("03","启动序列",             "7-Phase Initialization Flow"),
    ("04","一次请求的旅程",       "QueryLoop Execution Engine"),
    ("05","工具执行引擎",         "45+ Tools Concurrent Scheduling"),
    ("06","权限决策引擎",         "5-Mode Security Gate"),
    ("07","上下文压缩系统",       "6-Strategy Token Management"),
    ("08","记忆与提示组装",       "4-Layer Memory · Prompt Cache"),
    ("09","扩展能力三件套",       "Hooks · MCP · Skills"),
    ("10","关键设计哲学",         "Core Insights & Summary"),
]
for i, (num, zh, en) in enumerate(items):
    col, row = i % 2, i // 2
    x0 = 0.5 + col * 6.55
    y0 = 1.2 + row * 1.18
    add_box(sld, x0, y0, 0.52, 0.52, num, 14, bold=True,
            fill=BLUE, text_color=WHITE, rounded=True)
    add_txt(sld, x0+0.65, y0+0.02, 5.5, 0.35, zh, 15, bold=True, color=DARK)
    add_txt(sld, x0+0.65, y0+0.38, 5.5, 0.3, en, 11, color=GRAY, font=EN)
print("✓ 02 目录")

# ══════════════════════════════════════════════════════════
# 幻灯片 03 — 这台机器是什么
# ══════════════════════════════════════════════════════════
sld = prs.slides.add_slide(blank)
add_rect(sld, 0, 0, W, H, fill=WHITE, rounded=False)
slide_header(sld, "01  |  这台机器是什么？", "Harness Engineering Paradigm")

# 左：误解
add_box(sld, 0.4, 1.25, 3.9, 0.48, "❌  常见误解", 14, bold=True,
        fill=LRED, text_color=RED, rounded=True)
shp = add_rect(sld, 0.4, 1.82, 3.9, 1.45,
               fill=RGBColor(0xff,0xf5,0xf5), line_color=RED, rounded=True)
tb = sld.shapes.add_textbox(Inches(0.55), Inches(1.95), Inches(3.6), Inches(1.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
run_add(p, "Claude Code 只是一个", 13, color=DARK)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
run_add(p2, "「在终端里问问题」", 15, bold=True, color=RED)
p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.LEFT
run_add(p3, "的聊天工具", 13, color=DARK)

add_arrow_right(sld, 4.45, 3.3, 0.45, 0.25, color=GRAY)

add_box(sld, 0.4, 3.65, 3.9, 0.48, "✅  真相", 14, bold=True,
        fill=LGREEN, text_color=GREEN, rounded=True)
add_box(sld, 0.4, 4.22, 3.9, 1.3,
        "它是一套精密的\nHarness 工程体系", 17,
        bold=True, fill=RGBColor(0xf0,0xfd,0xf4),
        text_color=DARK, line_color=GREEN, rounded=True)

# 右：概念卡片
add_rect(sld, 5.0, 1.2, 8.1, 5.65, fill=ACCENT, line_color=LBLUE, rounded=True)
add_txt(sld, 5.25, 1.35, 7.5, 0.4,
        "Harness 工程范式  Harness Engineering Pattern", 14, bold=True, color=BLUE)

concepts = [
    ("⚙",  "Harness（驱动层）",
     "为 LLM 提供可操作的执行环境，让推理结果能真实影响文件系统、Shell、IDE、MCP 服务。这是 Claude Code 的「躯壳」。"),
    ("🧠", "Model（推理核心）",
     "Claude API，只负责推理，完全无状态。不持有任何会话状态——所有状态都由 Harness 管理。"),
    ("🔗", "核心关系",
     "Model 是大脑，Harness 是神经系统+肌肉。没有 Harness，推理停留在对话框；没有 Model，Harness 不知道该做什么。"),
]
for j, (icon, title, body) in enumerate(concepts):
    y = 1.95 + j * 1.3
    add_txt(sld, 5.2, y, 0.4, 0.38, icon, 18, color=BLUE)
    add_txt(sld, 5.7, y, 7.2, 0.35, title, 14, bold=True, color=NAVY)
    add_txt(sld, 5.7, y + 0.38, 7.1, 0.72, body, 12, color=DARK)

add_box(sld, 5.0, 6.62, 8.1, 0.52,
        "Claude Code  =  Claude API（推理） + Harness（执行环境）",
        13, bold=True, fill=NAVY, text_color=WHITE, font=EN)
print("✓ 03 这台机器是什么")

# ══════════════════════════════════════════════════════════
# 幻灯片 04 — 四层体系
# ══════════════════════════════════════════════════════════
sld = prs.slides.add_slide(blank)
add_rect(sld, 0, 0, W, H, fill=WHITE, rounded=False)
slide_header(sld, "02  |  四层体系总览",
             "L1 Execution  ·  L2 Context  ·  L3 Control  ·  L4 Extension")

layers = [
    ("L4","扩展层  Extension Layer",
     "MCP Servers  ·  Hooks  ·  Skills  ·  Custom Commands",
     "外部能力注入，无限可扩展", LPURPL, PURPLE),
    ("L3","控制层  Control Layer",
     "Permissions  ·  Plan Mode  ·  Budget  ·  Security",
     "安全边界与行为约束", LORANG, ORANGE),
    ("L2","上下文层  Context Layer",
     "CLAUDE.md  ·  Memory Files  ·  Dialog History  ·  Prompt Cache",
     "知识注入与对话管理", LGREEN, GREEN),
    ("L1","执行层  Execution Layer",
     "QueryEngine  →  QueryLoop  →  StreamingToolExecutor",
     "唯一真正执行的层，驱动整个循环", SKYBLUE, BLUE),
]
widths = [13.0, 11.8, 10.6, 9.4]
bx = 0.17
by = 1.2
bh = 1.2
gap = 0.06

for i, (num, title, detail, note, bg, fg) in enumerate(layers):
    lw = widths[i]
    lx = bx + (13.0 - lw) / 2
    ly = by + i * (bh + gap)
    add_rect(sld, lx, ly, lw, bh, fill=bg, line_color=fg, rounded=True)
    add_box(sld, lx+0.12, ly+(bh-0.52)/2, 0.55, 0.52, num, 13, bold=True,
            fill=fg, text_color=WHITE, rounded=True)
    add_txt(sld, lx+0.82, ly+0.1, 7.5, 0.42, title, 15, bold=True, color=fg, font=EN)
    add_txt(sld, lx+0.82, ly+0.56, 6.5, 0.38, detail, 11, color=DARK, font=EN)
    add_txt(sld, lx+lw-3.6, ly+0.12, 3.3, 0.9, note, 12, color=GRAY,
            align=PP_ALIGN.RIGHT)

add_txt(sld, 0.4, 6.58, 12.5, 0.65,
        "关键认知：L1 是唯一真正「执行」的层。Claude API 在 L1 核心但完全无状态——"
        "所有状态（对话历史、权限规则、记忆文件）由 L2/L3 层的 Harness 持有。",
        11, color=GRAY)
print("✓ 04 四层体系")

# ══════════════════════════════════════════════════════════
# 幻灯片 05 — 启动序列
# ══════════════════════════════════════════════════════════
sld = prs.slides.add_slide(blank)
add_rect(sld, 0, 0, W, H, fill=WHITE, rounded=False)
slide_header(sld, "03  |  启动序列  ·  从命令行到 QueryLoop",
             "$ claude  →  7 Phases  →  QueryLoop")

phases = [
    ("P1","配置加载",   "三层覆盖：~/.claude/settings.json → .claude/settings.local.json → CLI flags",    LORANG, ORANGE),
    ("P2","上下文发现", "检测 git root；链式向上加载所有 CLAUDE.md；加载 Memory 文件（≤200行）",           LGREEN, GREEN),
    ("P3","扩展初始化", "预缓存 Skill 元数据（懒加载）；启动 MCP 子进程；注册 Hooks 事件处理器",          LPURPL, PURPLE),
    ("P4","工具注册",   "组装 tools[]：45+ 内置工具 + MCP 动态工具（Skill 工具此时不加载）",              LTEAL,  TEAL),
    ("P5","引擎实例化", "创建 QueryEngine 会话级单例，绑定 BudgetTracker、权限函数 useCanUseTool()",       SKYBLUE, BLUE),
    ("P6","提示组装",   "3路并行：systemPrompt + userContext + systemContext；预热 Prompt Cache",          LPINK,  PINK),
    ("P7","会话启动",   "Interactive：显示欢迎屏等待输入  |  Non-interactive：处理 stdin 直接进入循环",    LGREEN, GREEN),
]
col_cnt = 2
widths_p = [6.2, 6.2]
xs = [0.4, 6.75]
rows_per_col = [4, 3]

idx = 0
for ci in range(col_cnt):
    for ri in range(rows_per_col[ci]):
        num, zh, en, bg, fg = phases[idx]
        x = xs[ci]
        y = 1.2 + ri * 1.42
        add_rect(sld, x, y, widths_p[ci], 1.18, fill=bg, line_color=fg, rounded=True)
        add_box(sld, x+0.1, y+(1.18-0.52)/2, 0.52, 0.52, num, 13, bold=True,
                fill=fg, text_color=WHITE, rounded=True)
        add_txt(sld, x+0.75, y+0.08, widths_p[ci]-0.85, 0.38, zh, 14, bold=True, color=fg)
        add_txt(sld, x+0.75, y+0.52, widths_p[ci]-0.85, 0.6, en, 10, color=DARK, font=EN)
        if ri < rows_per_col[ci] - 1:
            add_arrow_down(sld, x+0.27, y+1.21, 0.18, 0.22, color=fg)
        idx += 1

# 右列末尾 → QueryLoop 标识
add_arrow_down(sld, xs[1]+0.27, 1.2+2*1.42+1.21, 0.18, 0.22, color=GREEN)
add_box(sld, xs[1], 1.2+3*1.42+0.08, 6.2, 0.48,
        "→  进入 QueryLoop（见 幻灯片 04）",
        13, bold=True, fill=NAVY, text_color=WHITE)
print("✓ 05 启动序列")

# ══════════════════════════════════════════════════════════
# 幻灯片 06 — QueryLoop
# ══════════════════════════════════════════════════════════
sld = prs.slides.add_slide(blank)
add_rect(sld, 0, 0, W, H, fill=WHITE, rounded=False)
slide_header(sld, "04  |  一次请求的旅程  ·  QueryLoop 执行引擎",
             "processUserInput → fetchSystemPromptParts → API → Tool? → Loop")

# 垂直流程（左侧）
flow = [
    ("①  用户输入",          "processUserInput()  斜杠命令预处理",            SKYBLUE, BLUE),
    ("②  提示组装",          "fetchSystemPromptParts()  3路并行 → API Payload",LGREEN,  GREEN),
    ("③  模型推理",          "调用 Claude API（无状态）→ 获取响应",            LORANG,  ORANGE),
    ("④  工具调用判断",      "stop_reason = tool_use → 进入工具执行分支",      LPURPL,  PURPLE),
    ("⑤  权限检查 + 执行工具","useCanUseTool() → StreamingToolExecutor",       LTEAL,   TEAL),
    ("⑥  输出/继续",         "end_turn → display_output + Stop Hooks + 等待",  LGREEN,  GREEN),
]
for i, (zh, en, bg, fg) in enumerate(flow):
    y = 1.2 + i * 0.95
    add_box(sld, 0.4, y, 5.5, 0.78, zh, 13, bold=True,
            fill=bg, text_color=fg, align=PP_ALIGN.CENTER)
    add_txt(sld, 6.05, y+0.17, 7.0, 0.42, en, 11, color=DARK, font=EN)
    if i < len(flow) - 1:
        add_arrow_down(sld, 0.4+2.5, y+0.8, 0.18, 0.18, color=fg)

# 循环箭头标注
add_rect(sld, 5.6, 1.2, 0.08, 5.4, fill=RGBColor(0xc7,0xd2,0xfe), rounded=False)
add_txt(sld, 5.72, 3.5, 0.5, 0.7, "循\n环", 12, bold=True,
        color=RGBColor(0x43,0x38,0xca))

# 右侧说明
add_rect(sld, 6.5, 1.2, 6.6, 5.5, fill=ACCENT, line_color=LBLUE, rounded=True)
add_txt(sld, 6.7, 1.38, 6.1, 0.38, "QueryLoop 三大核心数据结构", 13, bold=True, color=NAVY)

structures = [
    ("QueryState",      "跨轮次可变状态容器\n  turn_count · transition · continuation_count\n  has_attempted_compact",  BLUE),
    ("stop_reason",     "循环控制信号\n  end_turn  →  退出循环，输出响应\n  tool_use  →  继续，进入工具执行",             GREEN),
    ("TransitionReason","继续原因标记（debug 友好）\n  tool_result_continuation  |  compact_retry\n  max_tokens_recovery  |  transport_retry", ORANGE),
]
dy = 0.0
for title, body, fg in structures:
    add_txt(sld, 6.7, 1.95+dy, 6.1, 0.35, title, 13, bold=True, color=fg, font=EN)
    add_txt(sld, 6.7, 2.32+dy, 6.1, 0.75, body, 11, color=DARK, font=EN)
    dy += 1.28

add_box(sld, 0.4, 6.65, 12.5, 0.52,
        "关键认知：QueryLoop 是永动机，只有两个出口——stop_reason=end_turn（正常结束）或预算耗尽",
        12, bold=True, fill=NAVY, text_color=WHITE)
print("✓ 06 QueryLoop")

# ══════════════════════════════════════════════════════════
# 幻灯片 07 — 工具执行引擎
# ══════════════════════════════════════════════════════════
sld = prs.slides.add_slide(blank)
add_rect(sld, 0, 0, W, H, fill=WHITE, rounded=False)
slide_header(sld, "05  |  工具执行引擎  ·  45+ 工具的并发调度",
             "StreamingToolExecutor  ·  isConcurrencySafe  ·  persistedOutput")

# 并发安全工具
add_txt(sld, 0.4, 1.2, 5.8, 0.38, "并发安全（isConcurrencySafe = true）", 13, bold=True, color=GREEN)
for i, tool in enumerate(["Read","Glob","Grep","WebFetch","WebSearch"]):
    add_box(sld, 0.4+i*1.16, 1.68, 1.05, 0.52, tool, 12,
            fill=LGREEN, text_color=GREEN, font=EN)

# 串行执行工具
add_txt(sld, 0.4, 2.42, 5.8, 0.38, "串行执行（isConcurrencySafe = false）", 13, bold=True, color=RED)
for i, tool in enumerate(["Write","Edit","Bash","Agent","NotebookEdit"]):
    add_box(sld, 0.4+i*1.16, 2.9, 1.05, 0.52, tool, 12,
            fill=LRED, text_color=RED, font=EN)

# TrackedTool
add_txt(sld, 0.4, 3.65, 5.8, 0.38, "TrackedTool  —  工具执行状态封装", 13, bold=True, color=BLUE)
code = [
    ("interface TrackedTool {",              False, DARK),
    ("  id:                string",          False, DARK),
    ("  isConcurrencySafe: boolean  ← 并发标志", True, BLUE),
    ("  status: 'pending'|'running'|'done'|'error'", False, DARK),
    ("  abortController:  AbortController", False, DARK),
    ("}",                                   False, DARK),
]
for j, (line, hi, col) in enumerate(code):
    bg = ACCENT if hi else LGRAY
    add_box(sld, 0.4, 4.1+j*0.38, 5.8, 0.36, line, 11,
            fill=bg, text_color=col, align=PP_ALIGN.LEFT, font=EN, rounded=False)

# 右侧：输出持久化
add_rect(sld, 6.5, 1.2, 6.65, 5.55, fill=ACCENT, line_color=LBLUE, rounded=True)
add_txt(sld, 6.7, 1.35, 6.2, 0.38, "工具输出持久化策略  persistedOutput", 13, bold=True, color=NAVY, font=EN)

rows_p = [
    ("通用工具阈值",  "50 KB",                       BLUE),
    ("Bash 阈值",     "30 KB",                       ORANGE),
    ("存储路径",      ".task_outputs/tool-results/", TEAL),
    ("消息替换",      "<persisted-output> 预览标记",  PURPLE),
    ("后续读取",      "模型用 read_file 按需取完整内容",GREEN),
]
for j, (label, val, fg) in enumerate(rows_p):
    y = 1.95 + j * 0.9
    add_box(sld, 6.7, y, 2.4, 0.52, label, 12, fill=LGRAY, text_color=DARK)
    add_box(sld, 9.22, y, 3.7, 0.52, val, 12,
            fill=RGBColor(0xe8,0xf0,0xfe), text_color=fg, font=EN)

add_box(sld, 6.5, 6.52, 6.65, 0.52,
        "设计意图：避免大输出占满 Context，磁盘是无限的，Context 是有限的",
        11, fill=RGBColor(0x0d,0x25,0x55), text_color=RGBColor(0xba,0xd0,0xf8))
print("✓ 07 工具执行引擎")

# ══════════════════════════════════════════════════════════
# 幻灯片 08 — 权限决策引擎
# ══════════════════════════════════════════════════════════
sld = prs.slides.add_slide(blank)
add_rect(sld, 0, 0, W, H, fill=WHITE, rounded=False)
slide_header(sld, "06  |  权限决策引擎  ·  五模式安全门",
             "useCanUseTool()  ·  Deny > Allow > Confirm")

modes = [
    ("default",           "默认模式",    "Deny 规则优先 → Allow 规则 → 其余操作请求用户确认",   SKYBLUE,  BLUE,   "日常使用推荐"),
    ("acceptEdits",       "接受编辑",    "文件编辑类操作自动通过，破坏性命令（rm -rf 等）仍询问", LGREEN,   GREEN,  "开发过程提效"),
    ("plan",              "计划模式",    "只读——禁止所有写入和 Shell 执行，适合代码审查",        LORANG,   ORANGE, "安全审查场景"),
    ("dontAsk",           "免打扰",      "跳过所有交互确认，适合 CI/脚本自动化，风险较高",       LRED,     RED,    "CI / 脚本专用"),
    ("bypassPermissions", "绕过权限",    "最高特权：跳过所有安全检查，需显式配置才能启用",        LPURPL,   PURPLE, "仅限受信任环境"),
]
for i, (mode, zh, desc, bg, fg, scenario) in enumerate(modes):
    y = 1.2 + i * 1.04
    add_rect(sld, 0.4, y, 12.5, 0.92, fill=bg, line_color=fg, rounded=True)
    add_box(sld, 0.55, y+0.2, 2.55, 0.5, mode, 13, bold=True,
            fill=fg, text_color=WHITE, font=EN)
    add_txt(sld, 3.28, y+0.06, 1.6, 0.38, zh, 14, bold=True, color=fg)
    add_txt(sld, 3.28, y+0.5,  7.8, 0.38, desc, 12, color=DARK)
    add_box(sld, 11.35, y+0.22, 1.4, 0.48, scenario, 11,
            fill=fg, text_color=WHITE, rounded=True)

add_box(sld, 0.4, 6.55, 12.5, 0.55,
        "决策优先级：Deny 规则  ＞  Allow 规则  ＞  用户交互确认（Deny 永远最先判断）",
        13, bold=True, fill=NAVY, text_color=WHITE)
print("✓ 08 权限决策引擎")

# ══════════════════════════════════════════════════════════
# 幻灯片 09 — 上下文压缩
# ══════════════════════════════════════════════════════════
sld = prs.slides.add_slide(blank)
add_rect(sld, 0, 0, W, H, fill=WHITE, rounded=False)
slide_header(sld, "07  |  上下文压缩系统  ·  六策略 Token 管理",
             "autoCompact triggered at 70% usage")

# 为什么
add_box(sld, 0.4, 1.2, 4.2, 1.1,
        "为什么需要压缩？\n读1000行文件 ≈ 4000 tokens\n读30文件+20命令 = 100K+\n200K 窗口很快就满",
        12, fill=LRED, text_color=DARK, align=PP_ALIGN.LEFT, rounded=True)
add_arrow_right(sld, 4.75, 1.55, 0.4, 0.22, color=GRAY)
add_box(sld, 5.28, 1.2, 7.8, 1.1,
        "触发条件：上下文使用量达到 70% 自动触发（保留 30% 缓冲，不等到 100%）",
        13, fill=ACCENT, text_color=DARK, align=PP_ALIGN.LEFT, rounded=True)

# 六种策略
strategies = [
    ("0","persistedOutput",  "工具执行时",   "大输出(>50KB)写磁盘，消息替换为预览标记",  SKYBLUE, BLUE),
    ("1","microCompact",     "每轮开始",     "静默截断超旧的 tool_result 内容",          LGREEN,  GREEN),
    ("2","snipCompact",      "超阈值时",     "保留首尾，压缩中间对话历史",               LORANG,  ORANGE),
    ("3","reactiveCompact",  "实时检测",     "动态裁剪，优先保留最近内容",               LTEAL,   TEAL),
    ("4","contextCollapse",  "会话过长",     "多轮对话合并为压缩摘要",                   LPURPL,  PURPLE),
    ("5","autoCompact",      "70% 阈值",     "调用 Claude API 生成高质量摘要（最强）",   LRED,    RED),
]
cw = 6.38
for i, (num, name, when, desc, bg, fg) in enumerate(strategies):
    col, row = i % 2, i // 2
    x = 0.4 + col * (cw + 0.48)
    y = 2.5 + row * 1.55
    add_rect(sld, x, y, cw, 1.35, fill=bg, line_color=fg, rounded=True)
    add_box(sld, x+0.1, y+(1.35-0.5)/2, 0.48, 0.5, num, 13, bold=True,
            fill=fg, text_color=WHITE, rounded=True)
    add_txt(sld, x+0.7, y+0.08, cw-0.85, 0.38, name, 13, bold=True, color=fg, font=EN)
    add_txt(sld, x+0.7, y+0.52, cw-0.85, 0.3, f"触发：{when}", 10, color=GRAY)
    add_txt(sld, x+0.7, y+0.84, cw-0.85, 0.45, desc, 11, color=DARK)
print("✓ 09 上下文压缩")

# ══════════════════════════════════════════════════════════
# 幻灯片 10 — 记忆与提示组装
# ══════════════════════════════════════════════════════════
sld = prs.slides.add_slide(blank)
add_rect(sld, 0, 0, W, H, fill=WHITE, rounded=False)
slide_header(sld, "08  |  记忆与提示组装  ·  四层记忆体系",
             "CLAUDE.md → Dialog → Memory Files → Prompt Cache")

# 左：四层记忆
mem = [
    ("L1","CLAUDE.md",      "固定上下文，每次启动自动加载，可 git 提交，团队共享",SKYBLUE, BLUE),
    ("L2","对话历史",        "当前会话动态增长，70% 时触发压缩，会话结束后丢失",  LGREEN,  GREEN),
    ("L3","Memory Files",   ".claude/memory/MEMORY.md，跨会话持久，≤200行/25KB",LORANG,  ORANGE),
    ("L4","Prompt Cache",   "重复内容零费用，命中率 40-95%，节省 40-60% Token",  LPURPL,  PURPLE),
]
for i, (lnum, name, desc, bg, fg) in enumerate(mem):
    y = 1.2 + i * 1.28
    add_rect(sld, 0.4, y, 6.0, 1.12, fill=bg, line_color=fg, rounded=True)
    add_box(sld, 0.55, y+(1.12-0.52)/2, 0.55, 0.52, lnum, 12, bold=True,
            fill=fg, text_color=WHITE, rounded=True)
    add_txt(sld, 1.25, y+0.08, 4.9, 0.38, name, 14, bold=True, color=fg)
    add_txt(sld, 1.25, y+0.54, 4.9, 0.52, desc, 11, color=DARK)

# 右：fetchSystemPromptParts
add_rect(sld, 6.7, 1.2, 6.45, 5.6, fill=ACCENT, line_color=LBLUE, rounded=True)
add_txt(sld, 6.9, 1.35, 6.0, 0.38,
        "fetchSystemPromptParts()  —  3路并行", 13, bold=True, color=NAVY, font=EN)

parts_info = [
    ("systemPrompt",  "核心人格 + 工具描述 + 安全规则\n静态，Prompt Cache 命中率最高（60-95%）",   BLUE),
    ("userContext",   "CLAUDE.md + Memory Files + 用户偏好\n半静态，项目内基本稳定（80%+）",       GREEN),
    ("systemContext", "环境信息 + Hooks 指引 + 可用命令\n动态，每次对话略有不同",                  ORANGE),
]
for j, (pname, pdesc, fg) in enumerate(parts_info):
    y = 1.95 + j * 1.55
    add_box(sld, 6.9, y, 2.6, 0.65, pname, 13, bold=True,
            fill=LGRAY, text_color=fg, font=EN)
    add_txt(sld, 9.65, y+0.05, 3.25, 0.6, pdesc, 11, color=DARK, font=EN)
    if j < 2:
        add_txt(sld, 8.2, y+0.72, 0.5, 0.35, "⊕", 16, color=BLUE)

add_box(sld, 6.9, 6.6, 6.05, 0.38,
        "→  合并为最终 system prompt blocks[]",
        12, bold=True, fill=NAVY, text_color=WHITE, font=EN)

add_txt(sld, 0.4, 6.55, 6.0, 0.65,
        "Prompt Cache 关键：首次建立缓存，后续同目录请求命中，Token 成本降低 40-60%。",
        11, color=GRAY)
print("✓ 10 记忆与提示组装")

# ══════════════════════════════════════════════════════════
# 幻灯片 11 — 扩展三件套
# ══════════════════════════════════════════════════════════
sld = prs.slides.add_slide(blank)
add_rect(sld, 0, 0, W, H, fill=WHITE, rounded=False)
slide_header(sld, "09  |  扩展能力三件套  ·  Hooks · MCP · Skills",
             "Extension Layer: Event-driven  ·  Protocol-based  ·  Template-based")

three_ext = [
    ("Hooks",  "事件驱动的 Shell 扩展",
     ["PreToolUse  — 工具前，可阻止（exit 1）",
      "PostToolUse — 工具后，后处理",
      "Stop        — 最终输出时",
      "UserPromptSubmit — 输入预处理"],
     "exit 0 通过 · exit 1 阻止 · exit 2 注入消息",
     LORANG, ORANGE),
    ("MCP",    "Model Context Protocol\n外部能力注入协议",
     ["传输层：stdio / SSE / HTTP / WebSocket",
      "能力类型：Tools / Resources / Prompts",
      "JSON-RPC 2.0 标准协议",
      "连接失败非致命（warn + skip）"],
     "让 Claude 连接任意外部服务与工具",
     LTEAL, TEAL),
    ("Skills", "可复用任务执行模板",
     ["SKILL.md 定义触发词与执行流程",
      "元数据启动预缓存，内容懒加载",
      "斜杠命令 /skill-name 触发",
      "支持项目级 + 全局用户级"],
     "将复杂工作流沉淀为一行命令",
     LPURPL, PURPLE),
]
cw3 = 4.1
for i, (name, sub, items, summary, bg, fg) in enumerate(three_ext):
    x = 0.4 + i * (cw3 + 0.2)
    add_box(sld, x, 1.2, cw3, 0.65, name, 22, bold=True,
            fill=fg, text_color=WHITE, font=EN)
    add_txt(sld, x, 1.95, cw3, 0.5, sub, 12, color=GRAY, align=PP_ALIGN.CENTER)
    add_rect(sld, x, 2.55, cw3, 3.0, fill=bg, line_color=fg, rounded=True)
    for j, item in enumerate(items):
        add_txt(sld, x+0.18, 2.72+j*0.65, cw3-0.28, 0.58,
                f"• {item}", 11, color=DARK, font=EN)
    add_box(sld, x, 5.68, cw3, 0.65, summary, 11, bold=True,
            fill=fg, text_color=WHITE)

add_box(sld, 0.4, 6.58, 12.5, 0.62,
        "本质区别：Hooks 是事件响应（被动）；MCP 是能力注入（主动扩展工具集）；Skills 是工作流沉淀（复用模板）。三者正交，可同时使用。",
        12, fill=ACCENT, text_color=DARK, rounded=True)
print("✓ 11 扩展三件套")

# ══════════════════════════════════════════════════════════
# 幻灯片 12 — 关键设计哲学
# ══════════════════════════════════════════════════════════
sld = prs.slides.add_slide(blank)
add_rect(sld, 0, 0, W, H, fill=NAVY, rounded=False)
add_rect(sld, 0,    0, 0.1, H, fill=LBLUE, rounded=False)
add_rect(sld, W-0.1,0, 0.1, H, fill=LBLUE, rounded=False)

add_txt(sld, 0.4, 0.22, 12, 0.65,
        "10  |  关键设计哲学  ·  核心洞察与总结", 24, bold=True, color=WHITE)

insights = [
    ("状态分离",  "Claude API 完全无状态，所有状态由 Harness 管理。这保证了推理的纯粹性，也使 Harness 成为可扩展的工程平台。",  LBLUE),
    ("懒加载优先","Skill 内容不在启动时加载，工具输出超限才写磁盘。按需加载是 Context 高效利用的核心原则。",                   RGBColor(0x34,0xd3,0x99)),
    ("失败不阻断","MCP 连接失败只 warn 不中断会话；工具报错不退出 QueryLoop。局部故障不影响整体可用性。",                      RGBColor(0xfb,0xbf,0x24)),
    ("缓存即节省","Prompt Cache 对稳定内容命中率可达 95%，是降低 Token 成本最具性价比的手段。",                               RGBColor(0xc0,0x84,0xfc)),
    ("单例即一致","QueryEngine 是进程级单例，所有状态在一个对象内，保证会话的线性一致性，避免并发竞争。",                       RGBColor(0xf4,0x72,0xb6)),
]
cw5 = 5.9
for i, (title, body, fg) in enumerate(insights):
    col, row = i % 2, i // 2
    if i == 4:
        x = 0.4 + 6.55
        y = 1.1 + 2 * 2.5
    else:
        x = 0.4 + col * 6.55
        y = 1.1 + row * 2.5
    add_rect(sld, x, y, cw5, 2.25, fill=RGBColor(0x0d,0x20,0x42),
             line_color=fg, rounded=True)
    add_box(sld, x+0.12, y+0.1, cw5-0.22, 0.52, title, 16, bold=True,
            fill=fg, text_color=NAVY)
    add_txt(sld, x+0.22, y+0.75, cw5-0.4, 1.38, body, 12,
            color=RGBColor(0xca,0xd3,0xe3))

print("✓ 12 关键设计哲学")

# ══════════════════════════════════════════════════════════
# 保存
# ══════════════════════════════════════════════════════════
out = r"D:\Code\VibeBloging\InsightClaudeCode\ClaudeCode精密拆解.pptx"
prs.save(out)
print(f"\n✅  已保存：{out}")
